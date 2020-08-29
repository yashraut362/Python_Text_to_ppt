from pptx import Presentation
from pptx.util import Inches
import shutil

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = input("Enter Title Here For slide 1 ")
subtitle.text = input("Enter Subtitle Here For slide 1 ")


bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = input("Enter the Heading for slide 2 :- ")

tf = body_shape.text_frame
p = tf.add_paragraph()
str = open(input('Give Path for Text file or drop File here '), 'r').read()
p.text = str
p.level = 2


slide = prs.slides.add_slide(prs.slide_layouts[2])
table_placeholder = slide.shapes[1]
shapes = slide.shapes

shapes.title.text = 'Adding a Table'

rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

# set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# write column headings
table.cell(0, 0).text = 'Foo'
table.cell(0, 1).text = 'Bar'

# write body cells
table.cell(1, 0).text = 'Baz'
table.cell(1, 1).text = 'Qux'

blank_slide_layout = prs.slide_layouts[3]
slide = prs.slides.add_slide(blank_slide_layout)

# uncomment the below line if you want to give custom image.
#img_path = input('Enter The Image path or drop image here ')
#if you are using above line please comment the line given below.
img_path = 'input_img.jpg'

left = Inches(0.5)
height = Inches(5)
top = Inches(0.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)



prs.save('{slidename}.pptx'.format(slidename=input("Enter Name for PPT ")))
